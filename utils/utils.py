from functools import wraps
from flask import session, redirect, url_for, flash
import re
from collections import Counter
import hashlib


from io import BytesIO
import os

from db import get_categories
from io import BytesIO
import hashlib
import pdfkit
from werkzeug.utils import secure_filename
from pptx import Presentation
from psycopg2.extras import RealDictCursor
import psycopg2.errors
from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup
from psycopg2.extras import RealDictCursor

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
import re
STOP_WORDS = {
    "a", "ale", "aj", "je", "sú", "pre", "na", "do", "v", "z", "sa", "od", "u", "nie", "pod", "nad",
    "ktorý", "ktorá", "ktoré", "to", "tento", "táto", "tieto", "že", "ako", "čo", "či", "s", "o", "so",
    "obsahuje", "materiál", "prezentácia", "dokument", "úloha", "zadané"
}


KNOWN_PHRASES = [
    "5. ročník", "6. ročník", "9. ročník", "stredná škola", "základná škola",
    "výpočtová technika", "informatika", "dejepis", "prírodoveda", "cvičný test"
]

def generate_tags_from_text(text, max_tags=10):
    text = text.lower()


    matched_phrases = []
    for phrase in KNOWN_PHRASES:
        if phrase in text:
            matched_phrases.append(phrase)
            text = text.replace(phrase, "")  


    words = re.findall(r'\b\w{3,}\b', text)
    keywords = [w for w in words if w not in STOP_WORDS]


    combined = matched_phrases + keywords
    most_common = Counter(combined).most_common(max_tags)

    return [tag for tag, _ in most_common]



def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if "user_id" not in session:
            flash(_("Najprv sa prihláste.", "warning"))
            return redirect(url_for("login.login"))
        return f(*args, **kwargs)
    return decorated_function


def get_file_hash(file_storage):
    file_bytes = file_storage.read()  # prečíta celý obsah súboru
    file_storage.seek(0)  # vráti späť ukazovateľ na začiatok
    return hashlib.sha256(file_bytes).hexdigest()


def html_text_to_pptx(title: str, html_content: str, output_file: str = "presentation.pptx"):
    def set_slide_background(slide):

            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(173, 216, 230)

    soup = BeautifulSoup(html_content, "html.parser")
    prs = Presentation()

    # Úvodný slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Generovane ai"
    set_slide_background(title_slide)

    bullet_slide_layout = prs.slide_layouts[1]

    paragraphs = [p.get_text().strip() for p in soup.find_all("p") if p.get_text().strip()]
    
    current_title = None
    current_bullets = []

    slide_pattern = re.compile(r"^Slide\s*\d+:?\s*(.*)", re.IGNORECASE)

    for line in paragraphs:
        slide_match = slide_pattern.match(line)
        if slide_match:
            if current_title:
                slide = prs.slides.add_slide(bullet_slide_layout)
                slide.shapes.title.text = current_title
                content = slide.placeholders[1]
                if current_bullets:
                    content.text = current_bullets[0]
                    for bullet in current_bullets[1:]:
                        content.text_frame.add_paragraph().text = bullet
                set_slide_background(slide)
            current_title = slide_match.group(1).strip()
            current_bullets = []
        elif line.startswith("-"):
            current_bullets.append(line.lstrip("-").strip())

    if current_title:
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = current_title
        content = slide.placeholders[1]
        if current_bullets:
            content.text = current_bullets[0]
            for bullet in current_bullets[1:]:
                content.text_frame.add_paragraph().text = bullet
        set_slide_background(slide)
    prs.save(output_file)
    return output_file


import pdfkit

def html_text_to_pdf(html_content: str, output_path: str):
    """
    Vygeneruje PDF súbor z HTML reťazca s podporou UTF-8, základného formátovania
    a automatickým zalomením strany pred sekciou so správnymi odpoveďami.
    """
    import pdfkit
    from dotenv import load_dotenv
    load_dotenv()

    # Získaj cestu z .env
    wkhtmltopdf_path = os.getenv("WKHTMLTOPDF_PATH")

    if not wkhtmltopdf_path:
        raise EnvironmentError("WKHTMLTOPDF_PATH not set in .env file")

    # Konfiguruj pdfkit
    config = pdfkit.configuration(wkhtmltopdf=wkhtmltopdf_path)

    # ➔ Úprava obsahu: ak nájdeme "Správne odpovede", vložíme pred to page break
    if "Správne odpovede" in html_content or "Correct Answers" in html_content:
        html_content = re.sub(
            r"(Správne odpovede|Correct Answers)",
            r'<div style="page-break-before: always;"></div><h2>\1</h2>',
            html_content,
            flags=re.IGNORECASE
        )

    full_html = f"""
    <!DOCTYPE html>
    <html lang="sk">
    <head>
        <meta charset="UTF-8">
        <style>
            body {{
                font-family: DejaVu Sans, Arial, sans-serif;
                padding: 2rem;
                font-size: 14px;
                line-height: 1.6;
            }}
            h1, h2, h3, h4 {{
                color: #2c3e50;
            }}
            ul {{
                padding-left: 1.5rem;
            }}
            /* Špeciálne pre page breaky pri PDF */
            div[style*="page-break-before"] {{
                display: block;
                width: 100%;
                height: 1px;
                margin: 20px 0;
            }}
        </style>
    </head>
    <body>
        {html_content}
    </body>
    </html>
    """

    try:
        pdfkit.from_string(full_html, output_path, configuration=config)
    except Exception as e:
        raise RuntimeError(f"Chyba pri generovaní PDF: {e}")


ALLOWED_LEVELS = {'1','2','3'}
ALLOWED_TYPES = {"Poznámky", "Prezentácia", "Test", "Otázky a odpovede","Notes","Presentation","Test","Q&A"}


def validate_inputs(topic, level, material_type):
    errors = []
    
    if not topic or len(topic.strip()) < 3:
        errors.append(_("Téma musí mať aspoň 3 znaky."))

    if len(topic) > 100:
        errors.append(_("Téma je príliš dlhá (max. 100 znakov)."))
    if level not in ALLOWED_LEVELS:
        errors.append(_("Neplatná úroveň vzdelávania."))

    if material_type not in ALLOWED_TYPES:
        errors.append(_("Neplatný typ materiálu."))

    return errors

CATEGORY_NAME_TO_KEY = {
    "ZŠ": "zs",
    "SŠ": "ss",
    "VŠ": "vs"
}
from flask_babel import gettext as _
def load_categories_translated():
    raw_categories = get_categories()
    categories = []

    for cat in raw_categories:
        original_name = cat["name"]
        key = CATEGORY_NAME_TO_KEY.get(original_name)
        if key:
            label = _(key)  # preloží podľa jazyka
            categories.append({
                "id": cat["id"],
                "key": key,
                "label": label
            })
        else:
            # fallback - ak tam je niečo neznáme, zobrazí pôvodný názov
            categories.append({
                "id": cat["id"],
                "key": original_name,
                "label": original_name
            })
    
    return categories

from db import get_focuses, get_grades
from flask_babel import _

# Focusy - načítanie a preklad
def load_focuses_translated():
    raw_focuses = get_focuses()
    focuses = []

    for focus in raw_focuses:
        original_name = focus["name"]
        code = focus["code"]
        label = _(code)  # použijeme code na preklad

        focuses.append({
            "id": focus["id"],
            "key": code,
            "label": label
        })

    return focuses

# Grades (ročníky) - načítanie a preklad
def load_grades_translated():
    raw_grades = get_grades()
    grades = []

    for grade in raw_grades:
        original_name = grade["name"]
        code = grade["code"]
        label = _(code)  # použijeme code na preklad

        grades.append({
            "id": grade["id"],
            "key": code,
            "label": label
        })

    return grades


import pdfplumber
from pptx import Presentation

def extract_text_from_pdf(file_path):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"[ERROR] Failed to extract text from PDF: {e}")
    return text.strip()

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"[ERROR] Failed to extract text from PPTX: {e}")
    return text.strip()

import pdfplumber


import fitz  # PyMuPDF

def pdf_to_html_text(pdf_path: str) -> str:
    """
    Načíta PDF súbor a vygeneruje jednoduchý HTML reťazec pre editáciu.
    """
    doc = fitz.open(pdf_path)
    html_content = ""

    for page in doc:
        blocks = page.get_text("blocks")  # Získaj bloky textu
        blocks = sorted(blocks, key=lambda b: (b[1], b[0]))  # zoradenie zhora-dole, zľava-doprava

        for block in blocks:
            text = block[4].strip()
            if text:
                # Môžeme obaliť každý blok do <p>
                html_content += f"<p>{text}</p>\n"

        html_content += '<div style="page-break-after: always;"></div>\n'  # oddeliť strany

    return f"""<!DOCTYPE html>
<html lang="sk">
<head>
    <meta charset="UTF-8">
</head>
<body>
{html_content}
</body>
</html>"""



def pptx_to_html_text(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    output_lines = []
    slide_number = 1

    for idx, slide in enumerate(prs.slides):
        header = None
        bullets = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            # Špeciálne spracovanie pre prvý slide
            if idx == 0:
                # Ignoruj autora ("Generovane ai" v druhom placeholderi)
                if shape == slide.placeholders[1]:
                    continue
                else:
                    continue  # Ignoruj celý prvý slide (ak je len title + author)
            
            # Pre normálne slidy spracuj text
            lines = text.splitlines()

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                if header is None:
                    header = line  # prvý riadok v slide - použijeme ako header
                else:
                    bullets.append(line)

        # Spracuj obsah slidu (okrem prvého úvodného)
        if header:
            output_lines.append(f"<p>Slide {slide_number}: {header}</p>")
            slide_number += 1

        for bullet in bullets:
            if bullet.startswith(("-", "•", "–")):
                bullet_text = bullet.lstrip("-•–").strip()
            else:
                bullet_text = bullet.strip()
            output_lines.append(f"<p>- {bullet_text}</p>")

    return "\n".join(output_lines)



